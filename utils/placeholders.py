import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER

def analizar_presentacion():
    # Construir la ruta al archivo PPTX
    ruta_templates = os.path.join(os.getcwd(), "..", "templates")
    ppt_file = "Modelo reporte 1.pptx"
    ruta_ppt = os.path.join(ruta_templates, ppt_file)

    if not os.path.exists(ruta_ppt):
        print(f"Error: No se encontró el archivo en la ruta: {ruta_ppt}")
        return

    try:
        # Abrir la presentación
        prs = Presentation(ruta_ppt)
        
        print("\nAnálisis detallado de la presentación:")
        print(f"Total de diapositivas: {len(prs.slides)}")
        
        # Analizar cada diapositiva y su layout
        for idx, slide in enumerate(prs.slides, 1):
            print(f"\nAnálisis de Diapositiva {idx}:")
            
            # Información del layout
            print("\nInformación del Layout:")
            print(f"Nombre del layout: {slide.slide_layout.name}")
            
            # Analizar placeholders del layout
            print("\nPlaceholders definidos en el layout:")
            for ph in slide.slide_layout.placeholders:
                try:
                    print(f"  - Layout Placeholder:")
                    print(f"    Nombre: {ph.name}")
                    print(f"    ID: {ph.placeholder_format.idx}")
                    if hasattr(ph, 'text'):
                        print(f"    Texto actual: {ph.text if ph.text else '[Vacío]'}")
                except Exception as e:
                    print(f"    Error al analizar placeholder: {str(e)}")
            
            # Analizar todas las formas en la diapositiva
            print("\nFormas en la diapositiva:")
            for shape in slide.shapes:
                try:
                    print(f"\n  Forma:")
                    print(f"    Tipo: {shape.shape_type}")
                    print(f"    Nombre: {shape.name}")
                    
                    # Verificar si es un placeholder
                    if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                        print(f"    Es placeholder: Sí")
                        print(f"    Placeholder tipo: {shape.placeholder_format.type}")
                        print(f"    Placeholder ID: {shape.placeholder_format.idx}")
                    else:
                        print(f"    Es placeholder: No")
                    
                    # Verificar si tiene texto
                    if hasattr(shape, 'text'):
                        texto = shape.text[:100] + '...' if len(shape.text) > 100 else shape.text
                        print(f"    Texto: {texto}")
                        
                except Exception as e:
                    print(f"    Error al analizar forma: {str(e)}")

        print("\nNota: Para que una forma sea un placeholder válido, debe estar definida")
        print("como tal en el layout de la diapositiva y tener un ID de placeholder asignado.")

    except Exception as e:
        print(f"Error al procesar la presentación: {e}")
        import traceback
        print(traceback.format_exc())

if __name__ == "__main__":
    analizar_presentacion()