import re
import os
from pptx import Presentation
from datetime import datetime
from typing import Optional
import json
import logging
from openai import AsyncOpenAI

logger = logging.getLogger(__name__)

def parse_agent_output(compliance_output: str) -> dict:
    """
    Intenta parsear el compliance_output como JSON y retornarlo como diccionario.
    Si falla, registra el error y retorna {}.
    """
    try:
        data = json.loads(compliance_output)
        if isinstance(data, dict):
            return data
        else:
            logger.error("El JSON parseado no es un diccionario.")
    except Exception as e:
        logger.error("Error al parsear el JSON del compliance_output: %s", e)
    return {}


def extract_placeholders_from_template(template_path: str) -> set:
    prs = Presentation(template_path)
    placeholders = set()
    pattern = re.compile(r"\[[^\]]+\]")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        matches = pattern.findall(run.text)
                        if matches:
                            for m in matches:
                                placeholders.add(m.strip())  # Normaliza eliminando espacios
    return placeholders


# Añadir las constantes con las opciones válidas
VALID_SECTORS = {
    "Cross sector",
    "Banca y entidades de crédito",
    "Seguros y reaseguros",
    "Mercado de valores e inversiones",
    "Fintech y criptomonedas",
    "Software / hardware",
    "Plataformas y redes sociales",
    "Ciberseguridad y protección de datos",
    "IA y blockchain",
    "Energías renovables y combustibles fósiles",
    "Gestión de residuos y recivclaje",
    "Industria química y emisiones industriales",
    "Automotriz y aeronáutica",
    "Producción de maquinaria y equipos",
    "Industria textil y de consumo",
    "Industria farmacéutica y biotecnología",
    "Dispositivos médicos y hospitales",
    "Seguridad alimentaria y nutrición",
    "Obras públicas y urbanismo",
    "Empresas inmobiliarias",
    "Arquitectura e ingeniería civil",
    "Aerolíneas y transporte marítimo",
    "Empresas de transporte terrestre y ferroviario",
    "Servicios de paquetería y distribución",
    "Comercio electrónico y marketplaces",
    "Grandes cadenas y tiendas físicas",
    "Franquicias y distribución",
    "Hoteles y alojamientos turísticos",
    "Agencias de viajes y aerolíneas",
    "Restaurantes y gastronomía",
    "Producción agrícola y ganadera",
    "Agroindustria y exportación",
    "Seguridad alimentaria y trazabilidad"
}

VALID_CATEGORIES = {
    "Ley Orgánica",
    "Ley Ordinaria",
    "Directiva",
    "Reglamento",
    "Decisión",
    "Recomendación",
    "Dictámen",
    "Real decreto",
    "Decreto ley",
    "Directriz o guideline",
    "Resolución",
    "Evaluación",
    "FAQ",
    "Memoria",
    "Código de buenas prácticas",
    "Certificación internacional",
    "Programa"
}

async def resolve_placeholder_async(
    placeholder: str, 
    manual_values: dict, 
    compliance_output: str = None, 
    llm: Optional[AsyncOpenAI] = None
) -> str:
    normalized = placeholder.strip()  # Normaliza el placeholder
    if normalized in manual_values:
        return manual_values[normalized]
    
    if normalized == "[Fecha de reporte]":
        return datetime.now().strftime("%d/%m/%Y")
    
    if normalized == "[Nombre de sector]":
        prompt = (
            f"Basado en la siguiente información de compliance:\n{compliance_output}\n\n"
            "Identifica y devuelve ÚNICAMENTE el sector más relevante de la siguiente lista "
            f"(sin añadir texto adicional):\n{', '.join(VALID_SECTORS)}"
        )
    elif normalized == "[Categoría]":
        prompt = (
            f"Basado en la siguiente información de compliance:\n{compliance_output}\n\n"
            "Identifica y devuelve ÚNICAMENTE la categoría normativa más apropiada de la siguiente lista "
            f"(sin añadir texto adicional):\n{', '.join(VALID_CATEGORIES)}"
        )
    else:
        prompt = (
            f"Basado en la siguiente información de compliance:\n{compliance_output}\n\n"
            f"Genera el contenido adecuado para el placeholder '{normalized}' "
            "sin añadir texto adicional. Responde únicamente con el valor a utilizar."
        )
    
    if llm is not None:
        try:
            response = await llm.chat.completions.create(
                model="gpt-4",
                messages=[
                    {"role": "system", "content": "Eres un asistente experto en normativa y compliance. Responde únicamente con el valor solicitado, sin texto adicional."},
                    {"role": "user", "content": prompt}
                ]
            )
            generated_value = response.choices[0].message.content.strip()
            
            # Validar valores para campos específicos
            if normalized == "[Nombre de sector]" and generated_value not in VALID_SECTORS:
                logger.warning(f"Sector generado '{generated_value}' no válido, usando valor por defecto")
                generated_value = "Cross sector"
            elif normalized == "[Categoría]" and generated_value not in VALID_CATEGORIES:
                logger.warning(f"Categoría generada '{generated_value}' no válido, usando valor por defecto")
                generated_value = "Reglamento"
                
            logger.info("Placeholder %s generado: %s", normalized, generated_value)
            return generated_value
        except Exception as e:
            logger.error("Error generando el valor para %s: %s", normalized, e)
            return f"[Error generando {normalized}]"
    else:
        logger.error("No se proporcionó una instancia de LLM para generar el valor para %s", normalized)
        return f"[No LLM for {normalized}]"


async def prepare_complete_placeholders(manual_values: dict, compliance_output: str, template_path: str, llm: AsyncOpenAI) -> dict:
    """
    Extrae los placeholders del template y, para cada uno, utiliza el valor generado por la LLM.
    """
    extracted = extract_placeholders_from_template(template_path)
    complete = {}
    for ph in extracted:
        # Llama a la función asíncrona para generar cada valor.
        complete[ph] = await resolve_placeholder_async(ph, manual_values, compliance_output, llm)
    return complete



def generate_ppt_report(template_path: str, placeholders: dict) -> str:
    """
    Genera un reporte PPTX a partir de una plantilla y un diccionario de placeholders.
    Se realiza un reemplazo global del texto en cada shape con un text_frame.
    """
    
    prs = Presentation(template_path)
    
    # Para cada slide
    for slide in prs.slides:
        # Para cada shape en la diapositiva
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text_frame:
                continue

            # Obtiene el texto completo del shape
            original_text = shape.text
            new_text = original_text

            # Reemplaza cada placeholder en el texto completo
            for ph, replacement in placeholders.items():
                new_text = new_text.replace(ph.strip(), replacement)

            
            # Limpia el text_frame y añade el nuevo texto
            text_frame = shape.text_frame
            # Borrar todos los párrafos existentes
            while len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
                p.clear()
                text_frame._element.remove(p._element)
            # Añade un nuevo párrafo con el texto modificado
            p = text_frame.add_paragraph()
            p.text = new_text

    # Directorio de salida
    output_dir = os.path.join(os.getcwd(), "output")
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    file_name = f"compliance_report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    output_path = os.path.join(output_dir, file_name)
    prs.save(output_path)
    
    return output_path

